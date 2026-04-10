"""
assemble_pptx.py — Assemble compute_layer_pptx_slides/ PNGs into Compute Layer.pptx

Run: source .venv/bin/activate && python3 assemble_pptx.py

Slide order:
  01 anchor
  02 dag
  03a estate_full
  04 current_arch
  05 runai_cost
  03b estate_runai   (after 05 — audience already knows the $3K price tag)
  06 two_paths
  07 quota_structure
  08 quota_procedure
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

# ── Paths ────────────────────────────────────────────────────────────────────
BASE_DIR  = os.path.dirname(__file__)
PNG_DIR   = os.path.join(BASE_DIR, 'compute_layer_pptx_slides')
OUT_FILE  = os.path.join(BASE_DIR, 'Compute Layer.pptx')

# ── Slide size — 16:9, matching 2880×1620 px at 200 dpi ────────────────────
# 2880 px / 200 dpi = 14.4 inches; 1620 px / 200 dpi = 8.1 inches
SLIDE_W = Inches(14.4)
SLIDE_H = Inches(8.1)

SLIDES = [
    'slide_01_anchor.png',
    'slide_02_dag.png',
    'slide_03a_estate_full.png',
    'slide_04_current_arch.png',
    'slide_05_runai_cost.png',
    'slide_03b_estate_runai.png',
    'slide_06_two_paths.png',
    'slide_07_quota_structure.png',
    'slide_08_quota_procedure.png',
    'slide_09_bionova_components.png',
    'slide_10_bionova_flow.png',
]

def build_pptx():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]  # completely blank layout

    for png_name in SLIDES:
        png_path = os.path.join(PNG_DIR, png_name)
        if not os.path.exists(png_path):
            print(f"  WARNING: missing {png_path} — skipping")
            continue

        slide = prs.slides.add_slide(blank_layout)

        # Insert image full-bleed
        slide.shapes.add_picture(
            png_path,
            left=Emu(0),
            top=Emu(0),
            width=SLIDE_W,
            height=SLIDE_H,
        )
        print(f"  Added {png_name}")

    prs.save(OUT_FILE)
    print(f"\nSaved: {OUT_FILE}")

if __name__ == '__main__':
    build_pptx()
