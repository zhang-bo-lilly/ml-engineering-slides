"""
assemble_pptx.py — Assembles slide_01..05.png into a 16:9 PPTX.
Output: scientific_ai_coe.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

SLIDES = [f'slide_{str(i).zfill(2)}.png' for i in range(1, 6)]
OUT    = 'scientific_ai_coe.pptx'

prs = Presentation()
prs.slide_width  = Inches(16)
prs.slide_height = Inches(9)

blank_layout = prs.slide_layouts[6]  # completely blank

for png in SLIDES:
    if not os.path.exists(png):
        print(f'  WARNING: {png} not found — skipping')
        continue
    slide = prs.slides.add_slide(blank_layout)
    slide.shapes.add_picture(png, 0, 0, prs.slide_width, prs.slide_height)
    print(f'added {png}')

prs.save(OUT)
print(f'\nsaved {OUT}')
