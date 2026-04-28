import json
import os
import subprocess
import sys

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu

BASE_DIR = os.path.dirname(os.path.abspath(__file__))

W, H = 2880, 1620

C = {
    'bg':      (255, 255, 255),
    'dark':    (26,  26,  26),
    'red':     (211, 32,  48),
    'gray':    (51,  51,  51),
    'caption': (153, 153, 153),
    'blue':    (59,  130, 246),
    'emerald': (16,  185, 129),
    'white':   (255, 255, 255),
    'card':    (243, 244, 246),
}

BREADCRUMB = 'ML ENGINEERING  ·  COMPUTE LAYER UPDATE'

_fcache = {}

def _font(name, size):
    key = (name, size)
    if key in _fcache:
        return _fcache[key]
    paths = {
        'Arial-Black': [
            '/System/Library/Fonts/Supplemental/Arial Black.ttf',
            '/Library/Fonts/Arial Black.ttf',
        ],
        'Arial-Bold': [
            '/System/Library/Fonts/Supplemental/Arial Bold.ttf',
            '/Library/Fonts/Arial Bold.ttf',
        ],
        'Arial': [
            '/System/Library/Fonts/Supplemental/Arial.ttf',
            '/Library/Fonts/Arial.ttf',
        ],
    }
    for p in paths.get(name, []):
        if os.path.exists(p):
            f = ImageFont.truetype(p, size)
            _fcache[key] = f
            return f
    f = ImageFont.load_default(size)
    _fcache[key] = f
    return f


def wrap_text(draw, text, fnt, max_w):
    words = text.split()
    lines, cur = [], []
    for word in words:
        test = ' '.join(cur + [word])
        if draw.textlength(test, font=fnt) > max_w and cur:
            lines.append(' '.join(cur))
            cur = [word]
        else:
            cur.append(word)
    if cur:
        lines.append(' '.join(cur))
    return lines


def line_height(draw, fnt, extra=8):
    bb = draw.textbbox((0, 0), 'Ag', font=fnt)
    return bb[3] - bb[1] + extra


def new_canvas():
    img = Image.new('RGB', (W, H), C['bg'])
    return img, ImageDraw.Draw(img)


def _draw_tracked(draw, x, y, text, font, fill, spacing=3):
    """Draw text with per-character letter spacing (PIL has no native tracking)."""
    for ch in text:
        draw.text((x, y), ch, font=font, fill=fill, anchor='lm')
        bb = draw.textbbox((0, 0), ch, font=font)
        x += bb[2] - bb[0] + spacing


def chrome(draw, page_num=None):
    _draw_tracked(draw, 194, 139, BREADCRUMB, _font('Arial-Bold', 22), C['red'], spacing=3)
    if page_num is not None:
        draw.text((W - 80, H - 35), f'{page_num:02d}',
                  font=_font('Arial', 28), fill=C['caption'], anchor='rm')
    draw.rectangle([(0, H - 18), (W, H)], fill=C['red'])


def hero(draw, line1, line2, y=189, size=98):
    max_w = W - 2 * 120
    while size > 40:
        f = _font('Arial-Black', size)
        if max(draw.textlength(line1, font=f), draw.textlength(line2, font=f)) <= max_w:
            break
        size -= 2
    f = _font('Arial-Black', size)
    draw.text((194, y), line1, font=f, fill=C['dark'], anchor='lt')
    bb1 = draw.textbbox((0, 0), line1, font=f)
    lh = bb1[3] - bb1[1] + 7
    draw.text((194, y + lh), line2, font=f, fill=C['red'], anchor='lt')
    bb2 = draw.textbbox((0, 0), line2, font=f)
    return y + lh + (bb2[3] - bb2[1])


def save_slide(img, name):
    p = os.path.join(BASE_DIR, name)
    img.save(p)
    print(f'  -> {p}')
    return p


def fmt_hours(h):
    if h >= 1000:
        return f'{h / 1000:.0f}K'
    return f'{h:.0f}'


def load_metrics():
    path = os.path.join(BASE_DIR, 'metrics.json')
    if not os.path.exists(path):
        raise FileNotFoundError(
            'metrics.json not found — run extract_metrics.py first'
        )
    with open(path) as f:
        return json.load(f)


# ── Slide 01 ──────────────────────────────────────────────────────────────

def slide_01(metrics):
    img, draw = new_canvas()
    chrome(draw, page_num=5)

    margin = 194
    gap = 40

    # Hero — slide-01 text at top
    hero(draw, 'LillyPod,', 'by the numbers.')

    # Stat tiles
    stats = [
        (str(metrics['projects']),    'active projects'),
        (str(metrics['users']),       'registered users'),
        (metrics['workloads_label'],  'workloads completed'),
        (metrics['cloud_cost_label'] + '*', 'cloud cost avoidance'),
    ]
    tile_w = (W - 2 * margin - 3 * gap) // 4
    tile_h = 210
    tile_y = 420

    for i, (num, label) in enumerate(stats):
        tx = margin + i * (tile_w + gap)
        draw.rectangle([(tx, tile_y), (tx + tile_w, tile_y + tile_h)], fill=C['card'])
        draw.rectangle([(tx, tile_y), (tx + tile_w, tile_y + 5)], fill=C['red'])

        fn = _font('Arial-Black', 96)
        nb = draw.textbbox((0, 0), num, font=fn)
        nw, nh = nb[2] - nb[0], nb[3] - nb[1]
        num_color = C['red'] if i == 3 else C['dark']
        draw.text(
            (tx + (tile_w - nw) // 2 - nb[0], tile_y + (tile_h - nh) // 2 - 18 - nb[1]),
            num, font=fn, fill=num_color,
        )

        fl = _font('Arial', 28)
        lb = draw.textbbox((0, 0), label, font=fl)
        draw.text(
            (tx + (tile_w - (lb[2] - lb[0])) // 2, tile_y + tile_h - 36),
            label, font=fl, fill=C['caption'],
        )

    # Projects bar chart — full width, natural height
    chart = Image.open(os.path.join(BASE_DIR, 'projects_gpuhours.png')).convert('RGB')
    orig_w, orig_h = chart.size
    chart_y = tile_y + tile_h + 36
    cw = W - 2 * margin
    ch = int(cw * orig_h / orig_w)
    img.paste(chart.resize((cw, ch), Image.LANCZOS), (margin, chart_y))

    # Caption anchored directly below the chart
    fc = _font('Arial', 32)
    draw.text((W // 2, chart_y + ch + 16),
              'One cluster, touching every part of the business.',
              font=fc, fill=C['caption'], anchor='mt')

    # Footnote — derive "as of" date from last entry in jobs_by_date
    as_of_iso = max(metrics['jobs_by_date'].keys())
    from datetime import date
    as_of = date.fromisoformat(as_of_iso).strftime('%B %-d, %Y')
    ff = _font('Arial', 24)
    draw.text((margin, H - 30), f'* Data as of {as_of}',
              font=ff, fill=C['caption'], anchor='lm')

    return save_slide(img, 'slide_01.png')


# ── Slide 02 ──────────────────────────────────────────────────────────────

CTAS = [
    {
        'title': 'Not on LillyPod yet?',
        'body': (
            "If you're in AIR but not yet running workloads, now is the time to think about "
            "where GPU compute fits in your roadmap. Getting started is simpler than you think "
            "— reach out and we'll set you up."
        ),
    },
    {
        'title': "Onboarded but haven't run workloads yet?",
        'body': (
            'Active workloads help keep the cluster healthy for everyone. '
            'The more of us running jobs, the more we cover each other. '
            'The cluster is ready — we\'re here to help you get your first job running.'
        ),
    },
    {
        'title': 'Bring your partners along.',
        'body': (
            'Working with teams outside AIR who could benefit from GPU compute? '
            'Introduce them to LillyPod. Every new user adds coverage '
            '— helping them is helping all of us.'
        ),
    },
]


def slide_02():
    img, draw = new_canvas()
    chrome(draw, page_num=6)
    hero_bottom = hero(draw, 'Not everyone has access like this.', 'At its best when we all use it.')

    margin = 194
    cta_top = hero_bottom + 30
    card_gap = 44
    card_w = W - 2 * margin
    pad_v = 80
    pad_h = 30

    fn = _font('Arial-Bold', 26)
    ft = _font('Arial-Black', 52)
    fb = _font('Arial', 34)
    lh_title = line_height(draw, ft, extra=10)
    lh_body = line_height(draw, fb, extra=8)
    text_w = card_w - 2 * pad_h

    cy = cta_top
    for i, cta in enumerate(CTAS):
        num_bb = draw.textbbox((0, 0), f'0{i + 1}', font=fn)
        num_h = num_bb[3] - num_bb[1]
        title_lines = wrap_text(draw, cta['title'], ft, text_w)
        body_lines = wrap_text(draw, cta['body'], fb, text_w)

        card_h = (5 + pad_v
                  + num_h + 8
                  + len(title_lines) * lh_title
                  + 14
                  + len(body_lines) * lh_body
                  + pad_v)

        draw.rectangle([(margin, cy), (margin + card_w, cy + card_h)], fill=C['card'])
        draw.rectangle([(margin, cy), (margin + card_w, cy + 5)], fill=C['red'])

        iy = cy + 5 + pad_v
        draw.text((margin + pad_h, iy), f'0{i + 1}', font=fn, fill=C['red'])
        iy += num_h + 8

        for line in title_lines:
            draw.text((margin + pad_h, iy), line, font=ft, fill=C['dark'])
            iy += lh_title
        iy += 14

        for line in body_lines:
            draw.text((margin + pad_h, iy), line, font=fb, fill=C['gray'])
            iy += lh_body

        cy += card_h + card_gap

    return save_slide(img, 'slide_02.png')


# ── Slide 02a (alternative: colored-panel cards) ──────────────────────────

def slide_02a():
    img, draw = new_canvas()
    chrome(draw, page_num=6)
    hero(draw, 'Not everyone has access like this.', 'At its best when we all use it.')

    c_margin = 192
    c_gap    = 44
    card_top    = 472
    card_bottom = 1444
    card_w   = (W - 2 * c_margin - 2 * c_gap) // 3   # 802 px
    header_h = 168

    for i, cta in enumerate(CTAS):
        cx = c_margin + i * (card_w + c_gap)

        # Card background
        draw.rectangle([(cx, card_top), (cx + card_w, card_bottom)], fill=(249, 250, 251))

        # Emerald header panel (matches target)
        draw.rectangle([(cx, card_top), (cx + card_w, card_top + header_h)], fill=C['emerald'])

        # Watermark — 35% white blend, right-aligned, vertically centered in header
        wm = f'{i + 1:02d}'
        fwm = _font('Arial-Black', 170)
        wm_bb = draw.textbbox((0, 0), wm, font=fwm)
        wm_color = (
            int(0.35 * 255 + 0.65 * C['emerald'][0]),
            int(0.35 * 255 + 0.65 * C['emerald'][1]),
            int(0.35 * 255 + 0.65 * C['emerald'][2]),
        )
        wm_x = cx + card_w - 20 - (wm_bb[2] - wm_bb[0]) - wm_bb[0]
        wm_y = card_top + (header_h - (wm_bb[3] - wm_bb[1])) // 2 - wm_bb[1]
        draw.text((wm_x, wm_y), wm, font=fwm, fill=wm_color)

        # Section label
        fl = _font('Arial-Bold', 18)
        draw.text((cx + 24, card_top + 18), 'LILLYPOD', font=fl, fill=C['white'])

        # Card title (white Arial Black, wrapping)
        ft = _font('Arial-Black', 38)
        text_w = card_w - 48
        lh_t = line_height(draw, ft, extra=6)
        ty = card_top + 50
        for line in wrap_text(draw, cta['title'], ft, text_w):
            draw.text((cx + 24, ty), line, font=ft, fill=C['white'])
            ty += lh_t

        # Body text
        fb = _font('Arial', 30)
        lh_b = line_height(draw, fb, extra=8)
        by = card_top + header_h + 28
        for line in wrap_text(draw, cta['body'], fb, text_w):
            draw.text((cx + 24, by), line, font=fb, fill=C['gray'])
            by += lh_b

    return save_slide(img, 'slide_02a.png')


# ── PPTX assembly ─────────────────────────────────────────────────────────

def assemble_pptx(slide_paths, filename='LillyPod_Update_20260429.pptx'):
    prs = Presentation()
    prs.slide_width = Emu(12192000)    # 13.333"
    prs.slide_height = Emu(6858000)    # 7.5"
    blank = prs.slide_layouts[6]
    for path in slide_paths:
        sl = prs.slides.add_slide(blank)
        sl.shapes.add_picture(path, 0, 0, prs.slide_width, prs.slide_height)
    out = os.path.join(BASE_DIR, filename)
    prs.save(out)
    print(f'  -> {out}')
    return out


# ── Entry point ───────────────────────────────────────────────────────────

if __name__ == '__main__':
    print('Extracting metrics...')
    subprocess.run([sys.executable, os.path.join(BASE_DIR, 'extract_metrics.py')], check=True)

    print('Running diagram generators...')
    for script in ['diagram_projects.py']:
        subprocess.run([sys.executable, os.path.join(BASE_DIR, script)], check=True)

    metrics = load_metrics()

    print('Building slides...')
    slides = [
        slide_01(metrics),
        slide_02(),
    ]

    print('Building alternative...')
    slide_02a()

    print('Assembling pptx...')
    assemble_pptx(slides)
    print('Done.')
