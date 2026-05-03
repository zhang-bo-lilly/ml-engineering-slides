import os
import subprocess
import sys
from PIL import Image, ImageDraw, ImageFont

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
W, H = 2880, 1620
MARGIN = 194
BOTTOM_BAR = 18
BREADCRUMB = 'COMPUTE CONSOLIDATION PLAN'

_fcache = {}

def _font(name, size):
    key = (name, size)
    if key in _fcache:
        return _fcache[key]
    paths = {
        'Arial-Black': [
            '/System/Library/Fonts/Supplemental/Arial Black.ttf',
            '/Library/Fonts/Arial Black.ttf',
        ],
        'Arial-Bold': [
            '/System/Library/Fonts/Supplemental/Arial Bold.ttf',
            '/Library/Fonts/Arial Bold.ttf',
        ],
        'Arial': [
            '/System/Library/Fonts/Supplemental/Arial.ttf',
            '/Library/Fonts/Arial.ttf',
        ],
    }
    for p in paths.get(name, []):
        if os.path.exists(p):
            f = ImageFont.truetype(p, size)
            _fcache[key] = f
            return f
    f = ImageFont.load_default(size)
    _fcache[key] = f
    return f


def _chrome(draw):
    draw.text((MARGIN, 139), BREADCRUMB, font=_font('Arial-Bold', 28),
              fill='#d32030', anchor='lm')
    draw.rectangle([0, H - BOTTOM_BAR, W, H], fill='#d32030')


def _hero(draw, line1, line2=None, dark=False, font_size=110):
    fh = _font('Arial-Black', font_size)
    c1 = '#ffffff' if dark else '#1a1a1a'
    y = 189
    bb1 = draw.textbbox((0, 0), line1, font=fh)
    h1 = bb1[3] - bb1[1]
    draw.text((MARGIN, y), line1, font=fh, fill=c1)
    if line2 is None:
        return y + h1 + 7 + h1  # reserve two-line height for layout consistency
    y2 = y + h1 + 7
    bb2 = draw.textbbox((0, 0), line2, font=fh)
    h2 = bb2[3] - bb2[1]
    draw.text((MARGIN, y2), line2, font=fh, fill='#d32030')
    return y2 + h2


def _wrap(draw, text, font, max_w):
    words = text.split()
    lines, cur = [], []
    for w in words:
        test = ' '.join(cur + [w])
        if draw.textlength(test, font=font) <= max_w:
            cur.append(w)
        else:
            if cur:
                lines.append(' '.join(cur))
            cur = [w]
    if cur:
        lines.append(' '.join(cur))
    return lines


def _text_h(draw, text, font):
    bb = draw.textbbox((0, 0), text, font=font)
    return bb[3] - bb[1]


# ── SLIDE 01 — Title / thesis (dark anchor) ─────────────────────────────────

def slide_01():
    img = Image.new('RGB', (W, H), '#1a2035')
    draw = ImageDraw.Draw(img)
    _chrome(draw)
    hero_bottom = _hero(draw, 'LillyPod evolves into an enterprise',
                        'platform for scientific-AI workflows.',
                        dark=True, font_size=90)

    # Thesis paragraph — picks up from MagTrain
    fthesis = _font('Arial', 36)
    thesis = ("MagTrain's GPUs, simulation-class GPUs, and fat CPU nodes consolidate onto a "
              'single Weka fast storage fabric under Run:ai. md3 joins the shared namespace '
              'without forced migration. Future compute requests route to the platform rather '
              'than spawning new purpose-built environments.')
    thesis_lines = _wrap(draw, thesis, fthesis, W - 2 * MARGIN)
    ty = hero_bottom + 44
    for line in thesis_lines:
        draw.text((MARGIN, ty), line, font=fthesis, fill='#8899b8')
        ty += 52

    # Lead-in line
    flead = _font('Arial-Bold', 36)
    ty += 20
    draw.text((MARGIN, ty), 'Two compounding risks make this time-sensitive.',
              font=flead, fill='#ffffff')
    ty += _text_h(draw, 'M', flead) + 40

    content_top = ty
    content_bottom = H - BOTTOM_BAR - 48
    GAP = 48
    card_w = (W - 2 * MARGIN - GAP) // 2
    BORDER_L = 6

    risks = [
        ('RISK 1', 'False cost narrative',
         'Island environments on cheap GPUs appear cost-effective until they hit their ceiling. '
         'By then the narrative has already set. Countering it after the fact is significantly '
         'harder than making the case now.'),
        ('RISK 2', 'Structural intake',
         'Without a deliberate change in how compute requests are handled, the island pattern '
         'continues regardless of what the platform builds. Research IT is now within AIR — '
         'the change required is operational, not structural.'),
    ]

    fl = _font('Arial-Bold', 30)
    ft = _font('Arial-Black', 66)
    fb = _font('Arial', 40)

    # Size cards to content + fixed padding, then center the block vertically
    CARD_PAD = 44
    BODY_LEAD = 58
    text_w = card_w - BORDER_L - 28 - 52
    max_content_h = 0
    for _, title, body in risks:
        label_h = _text_h(draw, 'RISK 1', fl)
        title_h = _text_h(draw, title, ft)
        n_lines = len(_wrap(draw, body, fb, text_w))
        content_h = label_h + 16 + title_h + 28 + n_lines * BODY_LEAD
        max_content_h = max(max_content_h, content_h)

    card_h = max_content_h + CARD_PAD * 2
    card_top = content_top + 24

    for j, (label, title, body) in enumerate(risks):
        cx = MARGIN + j * (card_w + GAP)
        draw.rectangle([cx, card_top, cx + card_w, card_top + card_h],
                       fill='#243050', outline='#3a4a6a', width=1)
        draw.rectangle([cx, card_top, cx + BORDER_L, card_top + card_h],
                       fill='#d32030')

        text_x = cx + BORDER_L + 28
        ty_c = card_top + CARD_PAD
        draw.text((text_x, ty_c), label, font=fl, fill='#d32030')
        ty_c += _text_h(draw, label, fl) + 16
        draw.text((text_x, ty_c), title, font=ft, fill='#ffffff')
        ty_c += _text_h(draw, title, ft) + 28
        for line in _wrap(draw, body, fb, text_w):
            draw.text((text_x, ty_c), line, font=fb, fill='#8899b8')
            ty_c += BODY_LEAD

    out = os.path.join(BASE_DIR, 'slide_01.png')
    img.save(out)
    print(f'saved {out}')


# ── SLIDE 02 — Platform map ──────────────────────────────────────────────────

def slide_02():
    img = Image.new('RGB', (W, H), '#ffffff')
    draw = ImageDraw.Draw(img)
    _chrome(draw)
    hero_bottom = _hero(draw, 'One Weka fabric.', 'New hardware joins it.')

    diag = Image.open(os.path.join(BASE_DIR, 'diagram_platform_map.png'))
    dw, dh = diag.size
    diag_top = hero_bottom + 72
    avail_w = W - 2 * MARGIN
    scale = min(1.0, avail_w / dw)
    if scale < 1.0:
        diag = diag.resize((int(dw * scale), int(dh * scale)), Image.LANCZOS)
        dw, dh = diag.size
    img.paste(diag, (MARGIN, diag_top))

    out = os.path.join(BASE_DIR, 'slide_02.png')
    img.save(out)
    print(f'saved {out}')


# ── SLIDE 03 — Cost ──────────────────────────────────────────────────────────

def slide_03():
    img = Image.new('RGB', (W, H), '#ffffff')
    draw = ImageDraw.Draw(img)
    _chrome(draw)
    hero_bottom = _hero(draw, '2026 capital cost breakdown.')

    content_top    = hero_bottom + 80
    content_bottom = H - BOTTOM_BAR - 40
    content_w      = W - 2 * MARGIN

    ft_row = _font('Arial', 40)
    fa_row = _font('Arial-Black', 52)
    fb_sub = _font('Arial', 30)

    PAD_V      = 16
    PAD_X      = 14
    SUB_INDENT = 28
    SUB_LEAD   = 36

    lillypod_subs = [
        ('256× RTX 6000 Pro (32 servers)',                     '~$5.7M'),
        ('16× CPU fat nodes',                                  '~$1.3M'),
        ('Network + cables (MagTrain integration + new nodes)','~$650K'),
        ('Weka hot tier +2PB (~12 additional nodes)',          '~$3.0M'),
    ]

    rows = [
        ('md3 — 256 RTX 6000 Pro (32 servers; expands to 512 total)', '~$5.7M',   False, 'simple'),
        ('LillyPod',                                                   '~$10.65M', False, 'expanded'),
        ('Shared DC infrastructure (UPS, racks, PDU, CDU, fire suppression)', '~$7.1M', False, 'simple'),
        ('2026 total',                                                 '~$23.45M', True,  'total'),
    ]

    def _base_h(label, font, kind):
        if kind == 'expanded':
            return PAD_V + _text_h(draw, label, font) + 8 + len(lillypod_subs) * SUB_LEAD + PAD_V
        return _text_h(draw, label, font) + PAD_V * 2

    base_heights = [_base_h(l, fa_row if t else ft_row, k) for l, _, t, k in rows]
    slack = max(0, content_bottom - content_top - sum(base_heights)) // len(rows)
    heights = [h + slack for h in base_heights]

    ty = content_top
    for i, (label, amount, is_total, kind) in enumerate(rows):
        font  = fa_row if is_total else ft_row
        row_h = heights[i]
        bg    = '#fff1f2' if is_total else ('#f9fafb' if i % 2 == 0 else '#ffffff')

        draw.rectangle([MARGIN, ty, MARGIN + content_w, ty + row_h], fill=bg)
        if is_total:
            draw.rectangle([MARGIN, ty, MARGIN + content_w, ty + row_h],
                           outline='#d32030', width=2)
            draw.rectangle([MARGIN, ty, MARGIN + 8, ty + row_h], fill='#d32030')
        else:
            draw.line([MARGIN, ty + row_h, MARGIN + content_w, ty + row_h],
                      fill='#e5e7eb', width=1)

        lbl_color = '#d32030' if is_total else '#1a1a1a'
        lbl_x = MARGIN + (8 + PAD_X if is_total else PAD_X)

        if kind == 'expanded':
            # Vertically center the content block inside the (slack-padded) row
            content_h = _text_h(draw, label, font) + 8 + len(lillypod_subs) * SUB_LEAD
            block_top = ty + (row_h - content_h) // 2
            lbl_h = _text_h(draw, label, font)
            draw.text((lbl_x, block_top), label, font=font, fill=lbl_color)
            amt_bb = draw.textbbox((0, 0), amount, font=fa_row)
            amt_w  = amt_bb[2] - amt_bb[0]
            draw.text((MARGIN + content_w - PAD_X - amt_w, block_top),
                      amount, font=fa_row, fill=lbl_color)
            sub_y = block_top + lbl_h + 8
            for sub_label, sub_amt in lillypod_subs:
                dot_cy = sub_y + _text_h(draw, sub_label, fb_sub) // 2
                draw.ellipse([lbl_x + SUB_INDENT - 8, dot_cy - 4,
                              lbl_x + SUB_INDENT,     dot_cy + 4], fill='#9ca3af')
                text_x = lbl_x + SUB_INDENT + 10
                draw.text((text_x, sub_y), sub_label, font=fb_sub, fill='#6b7280')
                label_w = draw.textlength(sub_label, font=fb_sub)
                draw.text((text_x + label_w + 12, sub_y), sub_amt, font=fb_sub, fill='#374151')
                sub_y += SUB_LEAD
        else:
            lbl_h = _text_h(draw, label, font)
            draw.text((lbl_x, ty + (row_h - lbl_h) // 2), label, font=font, fill=lbl_color)
            amt_bb = draw.textbbox((0, 0), amount, font=fa_row)
            amt_w  = amt_bb[2] - amt_bb[0]
            amt_h  = amt_bb[3] - amt_bb[1]
            draw.text((MARGIN + content_w - PAD_X - amt_w, ty + (row_h - amt_h) // 2),
                      amount, font=fa_row, fill=lbl_color)

        ty += row_h

    out = os.path.join(BASE_DIR, 'slide_03.png')
    img.save(out)
    print(f'saved {out}')


# ── SLIDE 04 — Capability delta ──────────────────────────────────────────────

def slide_04():
    img = Image.new('RGB', (W, H), '#ffffff')
    draw = ImageDraw.Draw(img)
    _chrome(draw)
    hero_bottom = _hero(draw, 'Capability delta.', font_size=95)

    content_top = hero_bottom + 72
    content_bottom = H - BOTTOM_BAR - 40
    CALLOUT_H = 210
    panels_bottom = content_bottom - CALLOUT_H - 32

    GAP = 48
    col_w = (W - 2 * MARGIN - GAP) // 2
    left_x = MARGIN
    right_x = MARGIN + col_w + GAP
    panel_h = panels_bottom - content_top
    HEADER_H = 64
    PAD = 36

    fhdr_panel = _font('Arial-Bold', 26)
    fb_b = _font('Arial', 36)
    BODY_LEAD = 56

    # Panel backgrounds
    draw.rectangle([left_x, content_top, left_x + col_w, panels_bottom],
                   fill='#f9fafb', outline='#e5e7eb', width=1)
    draw.rectangle([right_x, content_top, right_x + col_w, panels_bottom],
                   fill='#f0fdf4', outline='#bbf7d0', width=1)

    # Panel headers
    draw.rectangle([left_x, content_top, left_x + col_w, content_top + HEADER_H],
                   fill='#374151')
    draw.rectangle([right_x, content_top, right_x + col_w, content_top + HEADER_H],
                   fill='#10b981')
    draw.text((left_x + PAD, content_top + 16), 'TODAY', font=fhdr_panel, fill='#ffffff')
    draw.text((right_x + PAD, content_top + 16), 'CONSOLIDATED PLATFORM',
              font=fhdr_panel, fill='#ffffff')

    text_w = col_w - 2 * PAD - 20

    today_items = [
        'Pipeline crosses cluster boundaries — simulation, featurization, and model retraining each live on a different cluster',
        'Inter-cluster data transfer is a blocking step — hundreds of GB of trajectory files, hours of wall-clock time, manual handoff at every stage',
    ]
    platform_items = [
        'Pipeline is automatable end-to-end — no manual staging, no inter-cluster copy, no lost time at boundaries',
        'Cross-cluster dependencies become a scheduler problem, not a data movement problem',
    ]

    def draw_bullets(items, base_x, text_color, dot_color):
        # Distribute bullets evenly in the panel below the header
        item_data = []
        for item in items:
            lines = _wrap(draw, item, fb_b, text_w - 22)
            item_data.append((lines, len(lines) * BODY_LEAD))

        avail = panel_h - HEADER_H - PAD
        total_content = sum(h for _, h in item_data)
        gap = (avail - total_content) / (len(item_data) + 1)
        iy = content_top + HEADER_H + PAD + gap

        for lines, item_h in item_data:
            draw.ellipse([base_x + PAD, iy + 14, base_x + PAD + 14, iy + 28],
                         fill=dot_color)
            for line in lines:
                draw.text((base_x + PAD + 22, iy), line, font=fb_b, fill=text_color)
                iy += BODY_LEAD
            iy += gap

    draw_bullets(today_items,    left_x,  '#374151', '#9ca3af')
    draw_bullets(platform_items, right_x, '#166534', '#10b981')

    # VS team callout
    callout_top = panels_bottom + 32
    fc = _font('Arial', 34)
    callout_text = ('Virtual screening team: running on cloud L40S and RTX 6000 GPUs at '
                    'on-demand rates, jobs hitting memory limits and splitting across GPUs. '
                    'H200 (141 GB HBM3e) and B300 (192 GB HBM3e) on the consolidated '
                    'platform eliminate the split — no extra GPUs, no extra cost.')
    callout_lines = _wrap(draw, callout_text, fc, W - 2 * MARGIN - 8 - 26 - 24)
    content_w = W - 2 * MARGIN
    draw.rectangle([MARGIN, callout_top, MARGIN + content_w, callout_top + CALLOUT_H],
                   fill='#fff1f2', outline='#fecdd3', width=1)
    draw.rectangle([MARGIN, callout_top, MARGIN + 8, callout_top + CALLOUT_H], fill='#d32030')
    total_text_h = len(callout_lines) * 52
    cty = callout_top + (CALLOUT_H - total_text_h) // 2
    for line in callout_lines:
        draw.text((MARGIN + 8 + 18, cty), line, font=fc, fill='#b91c1c')
        cty += 52

    out = os.path.join(BASE_DIR, 'slide_04.png')
    img.save(out)
    print(f'saved {out}')


# ── SLIDE 05 — Migration path ────────────────────────────────────────────────

def slide_05():
    img = Image.new('RGB', (W, H), '#ffffff')
    draw = ImageDraw.Draw(img)
    _chrome(draw)
    hero_bottom = _hero(draw, 'Timeline.')

    content_top = hero_bottom + 44
    content_bottom = H - BOTTOM_BAR - 40
    content_w = W - 2 * MARGIN

    fhdr = _font('Arial-Bold', 26)
    fn_step = _font('Arial-Bold', 28)
    ft_step = _font('Arial-Black', 40)
    fb_step = _font('Arial', 30)
    ft_timing = _font('Arial-Bold', 34)

    ty = content_top
    draw.text((MARGIN, ty), 'First capacity online Q3 2026', font=fhdr, fill='#9ca3af')
    ty += 40

    steps = [
        ('Step 0', 'Near-term',      'md3 RTX 6000 Pro deployment',        'md3 expands to 512 GPUs · Grid Engine',                      False),
        ('Step 1', 'Q3 2026',        'MagTrain Weka fabric integration',   'H100 / H200 / L40S join LillyPod under Run:ai',              True),
        ('Step 2', 'Q3 2026',        'CPU fat node deployment',            'Bioinformatics workloads on platform',                        False),
        ('Step 3', 'Q4 2026',        'RTX 6000 Pro deployment to LillyPod','End-to-end scientific-AI workflows on one fabric',            False),
        ('Step 4', 'Q4 2026',        'Weka +2PB expansion',               'Storage proportional to GPU capacity increase',               False),
        ('Step 5', 'Q1 2027',        'md3 Weka fabric membership',        'Data silo eliminated · md3 stays on Grid Engine',            False),
    ]

    CALLOUT_H = 110
    table_bottom = content_bottom - CALLOUT_H - 28
    row_h = (table_bottom - ty) // len(steps)

    # Column x-positions (all relative to MARGIN + 20 start)
    COL_X = MARGIN + 20
    W_STEP   = 140
    W_TIMING = 260
    W_TITLE  = 870
    x_timing  = COL_X + W_STEP + 16
    x_title   = x_timing + W_TIMING + 16
    x_outcome = x_title + W_TITLE + 16
    W_OUTCOME = (MARGIN + content_w - 20) - x_outcome

    for idx, (step, timing, title, outcome, highlight) in enumerate(steps):
        bg = '#fffbeb' if highlight else ('#f9fafb' if idx % 2 == 0 else '#ffffff')
        draw.rectangle([MARGIN, ty, MARGIN + content_w, ty + row_h], fill=bg)
        if highlight:
            draw.rectangle([MARGIN, ty, MARGIN + 5, ty + row_h], fill='#f59e0b')

        row_mid = ty + row_h // 2
        tim_color = '#d97706' if highlight else '#1a1a1a'

        step_h = _text_h(draw, step, fn_step)
        draw.text((COL_X, row_mid - step_h // 2), step, font=fn_step, fill='#9ca3af')

        tim_h = _text_h(draw, timing, ft_timing)
        draw.text((x_timing, row_mid - tim_h // 2), timing, font=ft_timing, fill=tim_color)

        tit_h = _text_h(draw, title, ft_step)
        draw.text((x_title, row_mid - tit_h // 2), title, font=ft_step, fill='#1a1a1a')

        out_lines = _wrap(draw, outcome, fb_step, W_OUTCOME)
        out_total_h = len(out_lines) * 42
        oy = row_mid - out_total_h // 2
        for line in out_lines:
            draw.text((x_outcome, oy), line, font=fb_step, fill='#6b7280')
            oy += 42

        draw.line([MARGIN, ty + row_h, MARGIN + content_w, ty + row_h],
                  fill='#e5e7eb', width=1)
        ty += row_h

    # Step 1 dependency callout
    callout_top = ty + 28
    fc = _font('Arial', 30)
    callout_text = ('Step 1 depends on inter-row cabling — validate physical constraints with '
                    'Greg Johnson + Jonathan Klinginsmith before committing the Q3 timeline.')
    callout_lines = _wrap(draw, callout_text, fc, content_w - 8 - 26 - 24)
    draw.rectangle([MARGIN, callout_top, MARGIN + content_w, callout_top + CALLOUT_H],
                   fill='#fffbeb', outline='#fde68a', width=1)
    draw.rectangle([MARGIN, callout_top, MARGIN + 8, callout_top + CALLOUT_H], fill='#f59e0b')
    total_text_h = len(callout_lines) * 44
    cty = callout_top + (CALLOUT_H - total_text_h) // 2
    for line in callout_lines:
        draw.text((MARGIN + 8 + 18, cty), line, font=fc, fill='#92400e')
        cty += 44

    out = os.path.join(BASE_DIR, 'slide_05.png')
    img.save(out)
    print(f'saved {out}')


# ── PPTX assembly ────────────────────────────────────────────────────────────

def assemble_pptx():
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    prs.slide_width = Pt(W * 72 / 96)
    prs.slide_height = Pt(H * 72 / 96)
    blank = prs.slide_layouts[6]

    for i in range(1, 6):
        slide = prs.slides.add_slide(blank)
        png = os.path.join(BASE_DIR, f'slide_0{i}.png')
        slide.shapes.add_picture(png, 0, 0, prs.slide_width, prs.slide_height)

    out = os.path.join(BASE_DIR, 'deck.pptx')
    prs.save(out)
    print(f'saved {out}')


if __name__ == '__main__':
    import argparse
    parser = argparse.ArgumentParser()
    parser.add_argument('--slide', type=int, default=0,
                        help='0=all+pptx, 1-5=single slide')
    args = parser.parse_args()

    for script in ['diagram_clusters.py', 'diagram_platform_map.py']:
        subprocess.run([sys.executable, os.path.join(BASE_DIR, script)], check=True)

    if args.slide == 0:
        slide_01(); slide_02(); slide_03(); slide_04(); slide_05()
        assemble_pptx()
    elif args.slide == 1:
        slide_01()
    elif args.slide == 2:
        slide_02()
    elif args.slide == 3:
        slide_03()
    elif args.slide == 4:
        slide_04()
    elif args.slide == 5:
        slide_05()
